#!/usr/bin/env python3
"""
Fase 6 — Output Generator.

Lê os resultados de um deal já processado (agent_runs + synthesis_runs) e
gera:
  output/{deal_id}_analise.xlsx  — 8 abas (7 do Prompt Master + Financial Analysis)
  output/{deal_id}_executivo.pptx — 5 slides, estilo CFO/IC/McKinsey

Uso:
    python scripts/generate_outputs.py --deal-id <uuid>

Sobre a entrega: até o conector do Microsoft Drive existir (Fase 5, hoje
bloqueada esperando acesso de TI), esses dois arquivos ficam disponíveis
como "Artifact" da própria execução do GitHub Actions — dá pra baixar
direto da página do run, sem precisar de mais nenhum acesso. Quando o
Drive estiver pronto, isso vira um upload automático além do artifact
(ou no lugar dele).
"""
import argparse
import json
import os
import sys
import urllib.request
import urllib.error
from datetime import datetime, timezone

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paleta — reaproveita as cores de marca da BHub (mesmo tom da arquitetura)
# ---------------------------------------------------------------------------
NAVY = "0F1727"
CORAL = "F25461"
SUCCESS = "2E7D4F"
WARN = "B6592A"
DANGER = "B03A3A"
TEXT_DARK = "141414"
TEXT_MUTED = "6A6A68"
BORDER_GRAY = "D8D8D6"
WHITE = "FFFFFF"

RECOMMENDATION_COLOR = {"GO": SUCCESS, "NO-GO": DANGER, "CONDITIONAL_GO": WARN}


# ---------------------------------------------------------------------------
# Supabase (mesma abordagem do run_agent.py)
# ---------------------------------------------------------------------------
def supabase_request(method: str, path: str) -> dict:
    url = os.environ["SUPABASE_URL"].rstrip("/") + "/rest/v1/" + path
    key = os.environ["SUPABASE_SERVICE_ROLE_KEY"]
    headers = {"apikey": key, "Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    req = urllib.request.Request(url, headers=headers, method=method)
    try:
        with urllib.request.urlopen(req) as resp:
            return json.loads(resp.read().decode())
    except urllib.error.HTTPError as e:
        print(f"Erro Supabase [{method} {path}]: {e.code} {e.read().decode()}", file=sys.stderr)
        raise


def fetch_deal_bundle(deal_id: str) -> dict:
    deals = supabase_request("GET", f"deals?id=eq.{deal_id}")
    if not deals:
        raise SystemExit(f"Deal {deal_id} não encontrado.")

    runs = supabase_request(
        "GET",
        f"agent_runs?deal_id=eq.{deal_id}&select=*,agent_versions(version,agent_id,agents(name))",
    )
    by_agent = {}
    for r in runs:
        name = r["agent_versions"]["agents"]["name"]
        by_agent[name] = r

    synthesis = supabase_request(
        "GET",
        f"synthesis_runs?deal_id=eq.{deal_id}&select=*,agent_versions(version)&order=created_at.desc&limit=1",
    )

    # deal_data mais recente — necessário pra acessar `raw_extracted`
    # direto (ex.: mini_dre, calculada 100% em código na extração, nunca
    # passou pelo output de nenhum agente de IA — ver run_agent.py).
    deal_data = supabase_request(
        "GET", f"deal_data?deal_id=eq.{deal_id}&order=created_at.desc&limit=1",
    )

    return {
        "deal": deals[0],
        "agents": by_agent,
        "synthesis": synthesis[0] if synthesis else None,
        "raw_extracted": (deal_data[0].get("raw_extracted") or {}) if deal_data else {},
    }


# ---------------------------------------------------------------------------
# EXCEL — 2 abas (achado real em 26/08: pedido explícito do Thiago — só
# "Executive Summary" e "Financial Analysis" importam, o resto (Complexity,
# Integration Risks, Operational Risks, Strategic Opinion, CFO Synthesis,
# Viabilidade Financeira, AI Audit Log) "não serve para nada" no Excel.
# As 2 abas continuam self-contained (não referenciam as removidas) — a
# única informação que se perderia com a remoção pura e simples era a
# "Fonte da Margem Bruta" (que morava só em Viabilidade Financeira), agora
# movida pra dentro de Financial Analysis.
# ---------------------------------------------------------------------------
def build_excel(bundle: dict, path: str):
    deal = bundle["deal"]
    agents = bundle["agents"]
    synth = bundle["synthesis"]
    raw_extracted = bundle.get("raw_extracted", {})

    wb = Workbook()
    header_font = Font(bold=True, color=WHITE, size=11)
    header_fill = PatternFill("solid", fgColor=NAVY)
    title_font = Font(bold=True, size=14, color=TEXT_DARK)
    thin_border = Border(*[Side(style="thin", color=BORDER_GRAY)] * 4)
    wrap = Alignment(wrap_text=True, vertical="top")

    def style_header_row(ws, row, ncols):
        for c in range(1, ncols + 1):
            cell = ws.cell(row=row, column=c)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(vertical="center")

    def autosize(ws, widths):
        for i, w in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(i)].width = w

    # 1. Executive Summary --------------------------------------------------
    ws = wb.active
    ws.title = "Executive Summary"
    ws["A1"] = deal["name"]
    ws["A1"].font = Font(bold=True, size=18, color=NAVY)
    ws.merge_cells("A1:B1")

    complexity = agents.get("complexity", {}).get("output", {}) or {}
    rec = synth.get("recommendation") if synth else None
    rows = [
        ("Deal", deal["name"]),
        ("Status", deal.get("status", "")),
        ("Complexidade", complexity.get("classificacao", "não avaliado")),
        ("Recomendação", rec or "não avaliado"),
        ("Confiança da recomendação", synth.get("confidence") if synth else None),
    ]
    r = 3
    for label, value in rows:
        ws.cell(row=r, column=1, value=label).font = Font(bold=True)
        ws.cell(row=r, column=2, value=value)
        r += 1

    r += 1
    ws.cell(row=r, column=1, value="Sumário executivo").font = title_font
    r += 1
    if synth:
        cell = ws.cell(row=r, column=1, value=synth.get("output", {}).get("sumario_executivo", ""))
        cell.alignment = wrap
        ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=4)
        ws.row_dimensions[r].height = 90
    r += 2

    def write_list_section(ws, r, title, items):
        ws.cell(row=r, column=1, value=title).font = title_font
        r += 1
        for item in items or []:
            ws.cell(row=r, column=1, value=f"• {item}")
            r += 1
        return r + 1

    if synth:
        out = synth.get("output", {})
        r = write_list_section(ws, r, "Principais riscos", out.get("principais_riscos"))
        r = write_list_section(ws, r, "Próximos passos / condições", out.get("condicoes"))
    autosize(ws, [28, 60, 20, 20])

    # 2. Financial Analysis — schema forense v2.0, + mini-DRE (26/08) e a
    #    fonte/motivo da Margem Bruta (movida de Viabilidade Financeira,
    #    que deixou de ter aba própria) ------------------------------------
    ws = wb.create_sheet("Financial Analysis")
    fin = agents.get("financial_analysis", {}).get("output", {}) or {}
    vf = agents.get("viabilidade_financeira", {}).get("output", {}) or {}
    row = 1

    # 0. Mini-DRE por período (determinística, 100% código — ver
    # dre_balancete_parser.py:montar_mini_dre) — só aparece quando a DRE
    # foi reconhecida nesse deal.
    mini_dre = next(iter((raw_extracted.get("mini_dre") or {}).values()), None)
    if mini_dre and mini_dre.get("linhas"):
        linhas_dre = mini_dre["linhas"]
        ws.cell(row=row, column=1, value="0. DRE por Período (R$)").font = title_font
        row += 1
        headers = ["Linha"] + [l["periodo"] for l in linhas_dre]
        for i, h in enumerate(headers, start=1):
            ws.cell(row=row, column=i, value=h)
        style_header_row(ws, row, len(headers))
        row += 1
        campos = [
            ("Receita Bruta", "receita_bruta"), ("(-) Deduções", "deducoes"),
            ("(=) Receita Líquida", "receita_liquida"), ("(-) Despesa com Pessoal", "despesa_pessoal"),
            ("(-) Custo Sistemas", "custo_sistemas"), ("(-) Outras Despesas", "outras_despesas"),
            ("(=) Resultado", "resultado"), ("Margem %", "margem_pct"),
        ]
        for label, campo in campos:
            ws.cell(row=row, column=1, value=label).font = Font(bold=True) if "=" in label else Font()
            for i, l in enumerate(linhas_dre, start=2):
                v = l.get(campo)
                ws.cell(row=row, column=i, value=f"{v}%" if campo == "margem_pct" and v is not None else v)
            row += 1
        ws.cell(row=row, column=1, value=f"Fonte: DRE ({mini_dre['fonte']}) — cálculo top-down, diferente do EBITDA Bridge abaixo (bottom-up).").font = Font(italic=True, size=9, color=TEXT_MUTED)
        row += 2

    # Fonte da Margem Bruta (achado real em 26/08, pedido explícito do
    # Thiago: "deve ficar claro se foi usado o formulário ou excel e por
    # quê" — antes só existia na aba Viabilidade Financeira, que deixou
    # de ter aba própria; preservado aqui, não perdido no corte).
    if vf.get("margem_bruta_fonte"):
        fonte_legivel = {
            "dre_fino": "DRE (categorização automática)", "dre_hierarquia": "DRE (extração por hierarquia)",
            "formulario": "Formulário",
        }.get(vf.get("margem_bruta_fonte"), vf.get("margem_bruta_fonte"))
        ws.cell(row=row, column=1, value="Fonte da Margem Bruta (Bloco A)").font = Font(bold=True)
        ws.cell(row=row, column=2, value=fonte_legivel)
        row += 1
        if vf.get("margem_bruta_motivo"):
            cell = ws.cell(row=row, column=1, value=f"Motivo: {vf['margem_bruta_motivo']}")
            cell.alignment = wrap
            cell.font = Font(italic=True, size=9, color=TEXT_MUTED)
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=6)
            row += 1
        row += 1

    ws.cell(row=row, column=1, value="1. EBITDA Bridge").font = title_font
    row += 1
    bridge = fin.get("ebitda_bridge") or {}
    bridge_rows = [
        ("Lucro Líquido", bridge.get("lucro_liquido")), ("Resultado Financeiro", bridge.get("resultado_financeiro")),
        ("Tributos sobre Lucro", bridge.get("tributos_sobre_lucro")), ("D&A", bridge.get("d_a")),
        ("EBITDA Reportado", bridge.get("ebitda_reportado")), ("Margem Reportada %", bridge.get("margem_reportada_pct")),
        ("EBITDA Ajustado", bridge.get("ebitda_ajustado")), ("Margem Ajustada %", bridge.get("margem_ajustada_pct")),
    ]
    for label, value in bridge_rows:
        ws.cell(row=row, column=1, value=label).font = Font(bold=True)
        ws.cell(row=row, column=2, value=value)
        row += 1
    ajustes = bridge.get("ajustes") or []
    if ajustes:
        row += 1
        ws.cell(row=row, column=1, value="Ajustes de não-recorrência").font = Font(bold=True)
        row += 1
        for aj in ajustes:
            ws.cell(row=row, column=1, value=aj.get("descricao"))
            ws.cell(row=row, column=2, value=aj.get("valor"))
            ws.cell(row=row, column=3, value=aj.get("justificativa"))
            row += 1

    row += 2
    ws.cell(row=row, column=1, value="2. Anomalias Detectadas").font = title_font
    row += 1
    headers = ["Conta", "Período", "Severidade", "Narrativa"]
    for i, h in enumerate(headers, start=1):
        ws.cell(row=row, column=i, value=h)
    style_header_row(ws, row, len(headers))
    row += 1
    for a in fin.get("anomalias", []):
        ws.cell(row=row, column=1, value=a.get("conta")); ws.cell(row=row, column=2, value=a.get("periodo"))
        ws.cell(row=row, column=3, value=a.get("severidade")); ws.cell(row=row, column=4, value=a.get("narrativa"))
        row += 1
    if not fin.get("anomalias"):
        ws.cell(row=row, column=1, value="Nenhuma anomalia detectada (ou sem série temporal suficiente para avaliar)")
        row += 1

    row += 2
    ws.cell(row=row, column=1, value="3. Red Flags").font = title_font
    row += 1
    headers = ["Severidade", "Título", "Detalhe", "Valor"]
    for i, h in enumerate(headers, start=1):
        ws.cell(row=row, column=i, value=h)
    style_header_row(ws, row, len(headers))
    row += 1
    for flag in fin.get("red_flags", []):
        ws.cell(row=row, column=1, value=flag.get("severidade")); ws.cell(row=row, column=2, value=flag.get("titulo"))
        ws.cell(row=row, column=3, value=flag.get("detalhe")); ws.cell(row=row, column=4, value=flag.get("valor"))
        row += 1

    row += 2
    ws.cell(row=row, column=1, value="4. Perguntas para Due Diligence").font = title_font
    row += 1
    headers = ["Prioridade", "Pergunta", "Contexto"]
    for i, h in enumerate(headers, start=1):
        ws.cell(row=row, column=i, value=h)
    style_header_row(ws, row, len(headers))
    row += 1
    for q in fin.get("perguntas_dd", []):
        ws.cell(row=row, column=1, value=q.get("prioridade")); ws.cell(row=row, column=2, value=q.get("pergunta"))
        ws.cell(row=row, column=3, value=q.get("contexto"))
        row += 1

    row += 2
    kpis = fin.get("kpis_valuation") or {}
    if kpis:
        ws.cell(row=row, column=1, value="5. KPIs para Valuation").font = title_font
        row += 1
        for k, v in kpis.items():
            ws.cell(row=row, column=1, value=k); ws.cell(row=row, column=2, value=v)
            row += 1

    limitacoes = fin.get("limitacoes_dados") or []
    if limitacoes:
        row += 2
        ws.cell(row=row, column=1, value="Limitações dos dados disponíveis").font = title_font
        row += 1
        for lim in limitacoes:
            cell = ws.cell(row=row, column=1, value=f"• {lim}")
            cell.alignment = wrap
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=6)
            row += 1

    autosize(ws, [26, 40, 16, 16, 12, 45])

    wb.save(path)
    print(f"Excel salvo: {path}")


# ---------------------------------------------------------------------------
# PPT — 5 slides, estilo CFO / Investment Committee
# ---------------------------------------------------------------------------
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                  color=TEXT_DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return box


def _add_bullets(slide, left, top, width, height, items, size=14, color=TEXT_DARK, space_after=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = f"—  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(color)
        run.font.name = "Calibri"
    return box


def _badge(slide, left, top, width, height, text, fill_hex, text_color=WHITE, size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    # shape.shadow.inherit = False não é suficiente pro LibreOffice — remove
    # o elemento de efeito diretamente na árvore XML da forma.
    sp = shape._element.spPr
    for tag in ("a:effectLst", "a:effectDag"):
        el = sp.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag.split(":")[1])
        if el is not None:
            sp.remove(el)
    from pptx.oxml.ns import qn
    effect_lst = sp.makeelement(qn("a:effectLst"), {})
    sp.append(effect_lst)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = _rgb(text_color)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def _blank_slide(prs, bg_hex=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg_hex)
    return slide


def _page_title(slide, text):
    _add_textbox(slide, MARGIN, Inches(0.45), Inches(11.5), Inches(0.65),
                 text, size=24, bold=True, color=NAVY)


def _footer(slide, deal_name, dark_bg=False):
    color = "6B7280" if not dark_bg else "6B7A99"
    _add_textbox(slide, MARGIN, Inches(7.12), Inches(6), Inches(0.3),
                 f"CONFIDENCIAL  ·  {deal_name}", size=9, color=color)


def _estimar_altura_bullets(items, largura_pol, size_pt, space_after_pt=8):
    """Estima a altura (em Inches/EMU) que uma lista de bullets vai
    ocupar — usada pra calcular a posição do PRÓXIMO bloco dinamicamente,
    em vez de um valor fixo que estoura quando o texto é mais longo que
    o previsto (bug real achado na QA visual em 26/08: 6 "próximos
    passos" longos empurravam o último item pra cima do rodapé)."""
    chars_por_linha = max(10, int((largura_pol * 72) / (size_pt * 0.50)))
    total_linhas = 0
    for item in items:
        texto = f"—  {item}"
        total_linhas += max(1, -(-len(texto) // chars_por_linha))  # ceil division
    altura_pt = total_linhas * size_pt * 1.22 + len(items) * space_after_pt
    return Inches(altura_pt / 72)


def _add_table(slide, left, top, width, row_h, headers, rows, header_fill=NAVY,
                col_widths=None, size=11, header_size=11, bold_last_row=True,
                highlight_rows=None):
    """Tabela simples pra dados tabulares (mini-DRE) — mais legível e
    mais compacta que bullets pra esse tipo de conteúdo, achado real em
    26/08 (pedido do Thiago pra ter uma mini-DRE visível, no estilo de
    teaser de M&A)."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = row_h * n_rows
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    highlight_rows = highlight_rows or set()

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(header_fill)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        run = p.runs[0]
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = _rgb(WHITE)
        run.font.name = "Calibri"
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r, row_vals in enumerate(rows, start=1):
        is_result_row = bold_last_row and r == n_rows - 1
        is_highlight = (r - 1) in highlight_rows
        for c, val in enumerate(row_vals):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("EEF1F6" if is_highlight else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            run = p.runs[0]
            run.font.size = Pt(size)
            run.font.bold = is_result_row or is_highlight
            run.font.color.rgb = _rgb(NAVY if (is_result_row or is_highlight) else TEXT_DARK)
            run.font.name = "Calibri"
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return gshape


def _fmt_mil(v) -> str:
    """R$ mil, formato brasileiro — compacto, como o teaser de referência."""
    if v is None:
        return "—"
    return f"{v/1000:,.0f}".replace(",", "X").replace(".", ",").replace("X", ".")


def build_pptx(bundle: dict, path: str):
    deal = bundle["deal"]
    agents = bundle["agents"]
    synth = bundle["synthesis"]
    raw_extracted = bundle.get("raw_extracted", {})
    complexity = agents.get("complexity", {}).get("output", {}) or {}
    opinion = agents.get("opinion", {}).get("output", {}) or {}
    integration = agents.get("integration_risks", {}).get("output", {}) or {}
    operational = agents.get("operational_risks", {}).get("output", {}) or {}
    financial = agents.get("financial_analysis", {}).get("output", {}) or {}
    synth_out = (synth or {}).get("output", {})
    recommendation = (synth or {}).get("recommendation", "N/D")
    deal_name = deal.get("name", "Deal")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---- Slide 1: Deal Overview (fundo navy, alto contraste) -------------
    s = _blank_slide(prs, bg_hex=NAVY)
    _add_textbox(s, MARGIN, Inches(2.3), Inches(12.1), Inches(1.0),
                 deal_name, size=38, bold=True, color=WHITE)
    _add_textbox(s, MARGIN, Inches(3.3), Inches(12.1), Inches(0.5),
                 "M&A Analysis Engine — Sumário Executivo", size=15, color="C9CDD6")

    complexity_label = complexity.get("classificacao", "não avaliado")
    stats = [
        ("Complexidade", complexity_label),
        ("Recomendação", recommendation),
        ("Confiança", f"{round((synth or {}).get('confidence', 0) * 100)}%" if synth else "n/d"),
    ]
    x = MARGIN
    for label, value in stats:
        _add_textbox(s, x, Inches(4.5), Inches(3.8), Inches(0.35), label.upper(), size=11, color="8B93A8")
        _add_textbox(s, x, Inches(4.85), Inches(3.8), Inches(0.6), str(value), size=20, bold=True, color=WHITE)
        x += Inches(4.0)
    _add_textbox(s, MARGIN, Inches(7.0), Inches(8), Inches(0.3), "ESTRITAMENTE CONFIDENCIAL", size=9, color="6B7A99")

    # ---- Slide 2: Complexity & Key Risks (fundido com o antigo slide 3,
    # 26/08 — sugestão do Thiago: Complexity sozinho tinha muito espaço
    # vazio, Key Risks já era denso; critérios viram tabela compacta em
    # vez de linhas espaçadas, sobra espaço pra "Integração" do lado) ---
    s = _blank_slide(prs)
    _page_title(s, "Complexity & Key Risks")
    _badge(s, MARGIN, Inches(1.15), Inches(3.0), Inches(0.42), complexity_label.upper(),
           fill_hex=DANGER if "Alta" in complexity_label else (WARN if "Média" in complexity_label else SUCCESS), size=11)

    criteria = complexity.get("criterios_acionados", [])
    integracao_top = Inches(3.0)
    if criteria:
        crit_rows = [[c.get("criterio", ""), c.get("detalhe", "")] for c in criteria[:5]]
        crit_row_h = Inches(0.26)
        _add_table(s, MARGIN, Inches(1.7), Inches(5.8), crit_row_h, ["Critério", "Nível"], crit_rows,
                   col_widths=[Inches(4.1), Inches(1.7)], size=10, header_size=10)
        integracao_top = Inches(1.7) + crit_row_h * (len(crit_rows) + 1) + Inches(0.3)

    _add_textbox(s, MARGIN, integracao_top, Inches(5.8), Inches(0.35), "INTEGRAÇÃO", size=12, bold=True, color=CORAL)
    _add_bullets(s, MARGIN, integracao_top + Inches(0.4), Inches(5.8), Inches(7.0) - (integracao_top + Inches(0.4)),
                 (integration.get("riscos_integracao") or ["Sem riscos de integração identificados."])[:4], size=11, space_after=8)

    _add_textbox(s, Inches(6.9), Inches(1.35), Inches(5.8), Inches(0.35), "OPERACIONAL", size=12, bold=True, color=CORAL)
    op_pct = operational.get("concentracao_receita_top10_pct")
    op_items = operational.get("riscos_operacionais") or ["Sem riscos operacionais identificados."]
    body_top = Inches(1.8)
    if op_pct is not None:
        level = operational.get("nivel_concentracao", "")
        _badge(s, Inches(6.9), Inches(0.85), Inches(3.2), Inches(0.4),
               f"Top 10: {op_pct}% ({level})",
               fill_hex=DANGER if level == "Crítico" else (WARN if level == "Médio" else SUCCESS), size=12)
    _add_bullets(s, Inches(6.9), body_top, Inches(5.8), Inches(5.0), op_items[:5], size=12, space_after=10)
    _footer(s, deal_name)

    # ---- Slide 3: CFO / Strategic Assessment -------------------------------
    s = _blank_slide(prs)
    _page_title(s, "CFO / Strategic Assessment")

    fin_result = financial.get("resultado", "n/d")
    _badge(s, MARGIN, Inches(1.2), Inches(2.8), Inches(0.44), f"Financeiro: {fin_result}",
           fill_hex=DANGER if fin_result == "Crítico" else (WARN if fin_result == "Atenção" else SUCCESS), size=12)

    # Mini-DRE (26/08, pedido do Thiago) — tabela de verdade, não texto
    # corrido. Só aparece quando a DRE foi reconhecida (`mini_dre`
    # calculada 100% em código na extração); cai pro EBITDA Bridge
    # clássico (bottom-up, texto) quando não há DRE nesse deal.
    mini_dre = next(iter((raw_extracted.get("mini_dre") or {}).values()), None)
    table_top = Inches(1.85)
    if mini_dre and mini_dre.get("linhas"):
        linhas_dre = mini_dre["linhas"]
        headers = ["R$ mil"] + [l["periodo"] for l in linhas_dre]
        campos = [
            ("Receita Bruta", "receita_bruta"), ("(-) Deduções", "deducoes"),
            ("(=) Receita Líquida", "receita_liquida"), ("(-) Despesa com Pessoal", "despesa_pessoal"),
            ("(-) Custo Sistemas", "custo_sistemas"), ("(-) Outras Despesas", "outras_despesas"),
            ("(=) Resultado", "resultado"),
        ]
        rows = []
        for label, campo in campos:
            valores = [_fmt_mil(l.get(campo)) for l in linhas_dre]
            rows.append([label] + valores)
        rows.append(["Margem %"] + [f"{l['margem_pct']}%" if l.get("margem_pct") is not None else "—" for l in linhas_dre])
        col_widths = [Inches(2.3)] + [Inches((5.8 - 2.3) / len(linhas_dre))] * len(linhas_dre)
        row_h = Inches(0.34)
        _add_table(s, MARGIN, table_top, Inches(5.8), row_h, headers, rows,
                   col_widths=col_widths, size=10.5, header_size=10.5, highlight_rows={2})
        # Posição da nota de rodapé calculada a partir da altura REAL da
        # tabela (linhas de dado + cabeçalho), não um valor fixo — bug
        # real achado na QA visual (26/08): um valor fixo colidia com o
        # título "EBITDA BRIDGE" logo abaixo quando a tabela tinha mais
        # linhas que o previsto.
        nota_top = table_top + row_h * (len(rows) + 1) + Inches(0.08)
        _add_textbox(s, MARGIN, nota_top, Inches(5.8), Inches(0.35),
                     f"Fonte: DRE ({mini_dre['fonte']}) — não confundir com o EBITDA Bridge ao lado (metodologia diferente).",
                     size=9, color=TEXT_MUTED)
        body_bottom = nota_top + Inches(0.45)
    else:
        body_bottom = table_top
        fin_lines = []
        bridge = financial.get("ebitda_bridge") or {}
        if bridge.get("ebitda_reportado") is not None:
            fin_lines.append(f"EBITDA: R$ {bridge['ebitda_reportado']:,.0f} ({bridge.get('margem_reportada_pct', '?')}%)".replace(",", "."))
        if bridge.get("ajustes"):
            fin_lines.append(f"EBITDA Ajustado: R$ {bridge.get('ebitda_ajustado', 0):,.0f} ({len(bridge['ajustes'])} ajuste(s))".replace(",", "."))
        if fin_lines:
            _add_bullets(s, MARGIN, Inches(1.9), Inches(5.8), Inches(2.5), fin_lines, size=12, space_after=10)

    # EBITDA Bridge (bottom-up) — tabela, não mais 1 linha de texto corrida
    # (achado real em 26/08, pedido do Thiago: "onde está a bridge?" — o
    # texto corrido escondia a estrutura do cálculo). Fica na coluna
    # direita, empilhado com Red Flags/Parecer — a coluna esquerda já
    # está ocupada pela mini-DRE, colocar as duas tabelas empilhadas na
    # mesma coluna estourava o slide.
    bridge = financial.get("ebitda_bridge") or {}

    # Red flags (só os 3 mais críticos, resto fica no Excel)
    _add_textbox(s, Inches(6.9), Inches(1.2), Inches(5.8), Inches(0.35), "RED FLAGS PRINCIPAIS", size=12, bold=True, color=CORAL)
    flags_lines = [f"({f.get('severidade')}) {f.get('titulo')}" for f in (financial.get("red_flags") or [])[:3]]
    bridge_top = Inches(2.9)
    if flags_lines:
        altura_flags = _estimar_altura_bullets(flags_lines, 5.8, 11.5, space_after_pt=8)
        _add_bullets(s, Inches(6.9), Inches(1.65), Inches(5.8), altura_flags + Inches(0.1), flags_lines, size=11.5, space_after=8)
        bridge_top = Inches(1.65) + altura_flags + Inches(0.3)

    parecer_top = bridge_top
    if bridge.get("ebitda_reportado") is not None:
        _add_textbox(s, Inches(6.9), bridge_top, Inches(5.8), Inches(0.3), "EBITDA BRIDGE (BOTTOM-UP)", size=11, bold=True, color=CORAL)
        bridge_campos = [
            ("Lucro Líquido", bridge.get("lucro_liquido")),
            ("(+) Resultado Financeiro", bridge.get("resultado_financeiro")),
            ("(+) Tributos sobre Lucro", bridge.get("tributos_sobre_lucro")),
            ("(+) D&A", bridge.get("d_a")),
            ("(=) EBITDA Reportado", bridge.get("ebitda_reportado")),
        ]
        bridge_rows = [[label, _fmt_mil(v)] for label, v in bridge_campos]
        bridge_rows.append(["Margem %", f"{bridge.get('margem_reportada_pct', '?')}%"])
        bridge_row_h = Inches(0.26)
        bridge_table_top = bridge_top + Inches(0.35)
        _add_table(s, Inches(6.9), bridge_table_top, Inches(3.8), bridge_row_h,
                   ["R$ mil", "Valor"], bridge_rows, col_widths=[Inches(2.5), Inches(1.3)],
                   size=9.5, header_size=9.5, highlight_rows={4})
        parecer_top = bridge_table_top + bridge_row_h * (len(bridge_rows) + 1) + Inches(0.25)

    parecer = opinion.get("parecer_do_time", "")
    if parecer:
        _add_textbox(s, Inches(6.9), parecer_top, Inches(5.8), Inches(0.3), "PARECER DO TIME", size=12, bold=True, color=CORAL)
        _add_textbox(s, Inches(6.9), parecer_top + Inches(0.35), Inches(5.8), Inches(7.0) - (parecer_top + Inches(0.35)), parecer, size=10.5, color=TEXT_DARK)
    _footer(s, deal_name)

    # ---- Slide 4: Recommendation & Next Steps ------------------------------
    s = _blank_slide(prs)
    _page_title(s, "Recommendation & Next Steps")
    _badge(s, MARGIN, Inches(1.15), Inches(3.6), Inches(0.6), recommendation,
           fill_hex=RECOMMENDATION_COLOR.get(recommendation, TEXT_MUTED), size=17)

    # ATENÇÃO (bug real corrigido em 26/08): `sumario_executivo` pode
    # passar de 200 palavras (um parágrafo corrido) — numa caixa de
    # altura fixa isso sempre estourava, sobrepondo "PRÓXIMOS PASSOS"
    # logo abaixo (slide "completamente quebrado", reportado pelo
    # Thiago). Correção estrutural, não só cosmética: usa
    # `principais_riscos` (lista de bullets curtos que o cfo_synthesis
    # já produz separadamente) em vez do parágrafo corrido — bullets têm
    # altura previsível, o parágrafo longo nunca teve. O texto completo
    # continua disponível no Excel (aba Executive Summary).
    riscos_curtos = synth_out.get("principais_riscos") or []
    top = Inches(2.05)
    LARGURA_BULLETS = 11.8  # polegadas, bate com Inches(11.8) usado abaixo
    LIMITE_INFERIOR = Inches(6.95)  # topo do rodapé — nada pode passar disso

    riscos_curtos = riscos_curtos[:4]
    if riscos_curtos:
        altura_riscos = _estimar_altura_bullets(riscos_curtos, LARGURA_BULLETS, 12.5, space_after_pt=8)
        _add_textbox(s, MARGIN, top, Inches(11.8), Inches(0.3), "PRINCIPAIS RISCOS", size=12, bold=True, color=CORAL)
        _add_bullets(s, MARGIN, top + Inches(0.4), Inches(11.8), altura_riscos + Inches(0.1), riscos_curtos, size=12.5, space_after=8)
        next_steps_top = top + Inches(0.4) + altura_riscos + Inches(0.35)
    else:
        next_steps_top = Inches(2.4)

    next_steps = (synth_out.get("condicoes") or ["Nenhuma condição adicional."])[:5]
    # Salvaguarda (bug real achado na QA visual em 26/08): mesmo com
    # posição dinâmica, uma lista longa pode ainda não caber até o
    # rodapé — reduz itens (nunca deixa texto vazar) até caber, ou até
    # sobrar só 3 (mínimo útil; o resto sempre está completo no Excel).
    while len(next_steps) > 3:
        altura_estim = _estimar_altura_bullets(next_steps, LARGURA_BULLETS, 12.5, space_after_pt=8)
        if next_steps_top + Inches(0.4) + altura_estim <= LIMITE_INFERIOR:
            break
        next_steps = next_steps[:-1]
    altura_next_steps = _estimar_altura_bullets(next_steps, LARGURA_BULLETS, 12.5, space_after_pt=8)

    _add_textbox(s, MARGIN, next_steps_top, Inches(11.8), Inches(0.3), "PRÓXIMOS PASSOS", size=12, bold=True, color=CORAL)
    _add_bullets(s, MARGIN, next_steps_top + Inches(0.4), Inches(11.8), altura_next_steps + Inches(0.1), next_steps, size=12.5, space_after=8)
    _footer(s, deal_name)

    prs.save(path)
    print(f"PPT salvo: {path}")


def upload_to_storage(local_path: str, storage_path: str, content_type: str):
    """Sobe um arquivo pro bucket privado deal-outputs, via service_role
    (o frontend nunca escreve aqui, só lê depois com signed URL)."""
    url = os.environ["SUPABASE_URL"].rstrip("/") + "/storage/v1/object/deal-outputs/" + storage_path
    key = os.environ["SUPABASE_SERVICE_ROLE_KEY"]
    with open(local_path, "rb") as f:
        data = f.read()
    headers = {
        "apikey": key,
        "Authorization": f"Bearer {key}",
        "Content-Type": content_type,
        "x-upsert": "true",  # permite reprocessar sem dar erro de arquivo já existente
    }
    req = urllib.request.Request(url, data=data, headers=headers, method="POST")
    with urllib.request.urlopen(req) as resp:
        print(f"Upload OK: {storage_path} ({resp.status})")


def record_output_link(deal_id: str, synthesis_run_id: str | None, excel_path: str, pptx_path: str):
    """Sobe os 2 arquivos pro Storage e grava as referências em public.outputs
    — é isso que o frontend usa para gerar o link de download real (signed
    URL), sem precisar passar pelo GitHub Actions."""
    if not os.environ.get("GITHUB_REPOSITORY"):
        print("Fora do GitHub Actions (teste local) — não sobe nem grava outputs.")
        return

    excel_storage_ref = f"{deal_id}/analise.xlsx"
    pptx_storage_ref = f"{deal_id}/executivo.pptx"
    upload_to_storage(
        excel_path, excel_storage_ref,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    upload_to_storage(
        pptx_path, pptx_storage_ref,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    body = {
        "deal_id": deal_id,
        "synthesis_run_id": synthesis_run_id,
        "excel_ref": excel_storage_ref,
        "ppt_ref": pptx_storage_ref,
    }
    url = os.environ["SUPABASE_URL"].rstrip("/") + "/rest/v1/outputs"
    key = os.environ["SUPABASE_SERVICE_ROLE_KEY"]
    headers = {"apikey": key, "Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    req = urllib.request.Request(url, data=json.dumps(body).encode(), headers=headers, method="POST")
    with urllib.request.urlopen(req) as resp:
        print(f"outputs gravado: {resp.status}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--deal-id", required=True)
    args = parser.parse_args()

    os.makedirs("output", exist_ok=True)
    bundle = fetch_deal_bundle(args.deal_id)

    excel_path = f"output/{args.deal_id}_analise.xlsx"
    pptx_path = f"output/{args.deal_id}_executivo.pptx"
    build_excel(bundle, excel_path)
    build_pptx(bundle, pptx_path)
    record_output_link(args.deal_id, (bundle.get("synthesis") or {}).get("id"), excel_path, pptx_path)
